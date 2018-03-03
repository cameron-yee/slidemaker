from __future__ import print_function
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from PyLyrics import *

#import argparse

#def lyrics(singer, song):
#	albums = PyLyrics.getAlbums(singer=singer)
#	for a in albums:
#		tracks = a.tracks()
#		for track in tracks:
#			print (track)
#			print (track.getLyrics())

# def groupLyrics(artist, song):
	# lyrics = PyLyrics.getLyrics(artist, song)
	# print(lyrics)
	# lyricsArray = lyrics.split("\n")
	
	# a = 0
	# for x in lyricsArray:
		# if(x == ""):
			# a += 1
	# numOfSlides = a + 1
	# print (a)
	# lyricSlides = [[] for i in range(numOfSlides)]
	
	# listNum = 0
	# for x in lyricsArray:
		# if(x != ""):
			# lyricSlides[listNum].append(x)
		# else:
			# listNum += 1
	# return lyricsArray
	
def groupLyrics(artist, song):
	lyrics = PyLyrics.getLyrics(artist, song)
	lyricsArray = lyrics.split("\n\n")

	return [x.split("\n") for x in lyricsArray]
	
	
# def create_ppt(song, groupOfLyrics):
	# prs = Presentation()
	# title_slide_layout = prs.slide_layouts[0]
	# titleSlide = prs.slides.add_slide(title_slide_layout)
	# songTitle= titleSlide.shapes.title
	
	# lyrics_slide_layout = prs.slide_layouts[1]
	
	# for lyrics in groupOfLyrics:
			# lyricSlide = prs.slides.add_slide(lyrics_slide_layout)
			# lyricPlaceholder = lyricSlide.placeholders[1]
			# lyricContainer = lyricPlaceholder.text_frame
			
			# textParagraph = lyricContainer.paragraphs[0]
			# textParagraph.alignment = PP_ALIGN.CENTER

			# run = textParagraph.add_run()
			# run.text =  "\n".join(lyrics)
			# #This lines don't work
			# run.size = Pt(42)
	
	# songTitle.text = song
	
	# prs.save('Slides/' + song + '.pptx')
	# #output = song + '.pptx'
	

#uses slideMaster custom template	for formatting
def create_ppt(song, groupOfLyrics):
	prs = Presentation('template.pptx')
	titleSlide = prs.slides[0]
	#LYRIC_SLIDE_LAYOUT = prs.slides[1].slide_layout
	LYRIC_SLIDE_LAYOUT = prs.slide_layouts[1]
	songTitle = titleSlide.shapes.title
	
	for lyrics in groupOfLyrics:
		lyricSlide = prs.slides.add_slide(LYRIC_SLIDE_LAYOUT)
		
		lyricPlaceholder = lyricSlide.placeholders[1]
		lyricContainer = lyricPlaceholder.text_frame
		
		textParagraph = lyricContainer.paragraphs[0]
		textParagraph.alignment = PP_ALIGN.CENTER
		
		run = textParagraph.add_run()
		run.text = "\n".join(lyrics)
			
	songTitle.text = song
	
		
	prs.save('Slides/' + song + '.pptx')
	
	

	
if __name__ == "__main__": 
	song = raw_input("Song Name: ")
	artist = raw_input("Artist Name: ")
	#output = PyLyrics.getLyrics("Rend Collective", "Build Your Kingdom Here")
	groupOfLyrics = groupLyrics(artist, song)
	create_ppt(song, groupOfLyrics)
	response = requests.get("")
	#print (output)
